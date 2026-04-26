"""
parser_agent.py
───────────────
Responsible for:
  1. Extracting raw text from uploaded lecture (.ppt/.pdf/.docx)
  2. Running spaCy NLP to find candidate concepts
  3. Calling LLM to build ordered concept dependency graph
  4. Writing concept_graph to the blackboard

Perceives:  file_path  (on blackboard)
Acts:       concept_graph, current_concept  (written to blackboard)
"""

import json
import re
from collections import Counter

import spacy
import ollama

from base_agent import BaseAgent


class ParserAgent(BaseAgent):

    def __init__(self, blackboard):
        super().__init__("ParserAgent", blackboard)

        # load spaCy English model for NLP
        print("  [ParserAgent] Loading spaCy model...")
        try:
            self.nlp = spacy.load("en_core_web_sm")
        except OSError:
            print("  [ParserAgent] spaCy model not found.")
            print("  Run: python -m spacy download en_core_web_sm")
            raise

    # ── PERCEIVE ─────────────────────────────────────────────
    def perceive(self):
        """Read the file path from the blackboard."""
        file_path = self.blackboard.read("file_path")
        print(f"  [ParserAgent] PERCEIVE → file path: {file_path}")
        self.blackboard.log_thinking(
            "ParserAgent",
            f"I see a lecture file at: {file_path}. I will extract its concepts."
        )
        return {"file_path": file_path}

    # ── REASON ───────────────────────────────────────────────
    def reason(self, perception: dict) -> list:
        """
        Full pipeline:
          raw text → slides → noun phrases → clean candidates → LLM concept graph
        """
        path = perception["file_path"]

        # ── Step 1: extract raw text ──────────────────────
        print("\n  [ParserAgent] REASON Step 1: Extracting text from file...")
        raw_text = self._extract_text(path)
        print(f"  [ParserAgent] Extracted {len(raw_text)} characters of text")
        self.blackboard.log_thinking(
            "ParserAgent",
            f"Extracted {len(raw_text)} characters. Splitting into slides..."
        )

        # ── Step 2: split into slides ─────────────────────
        slides = self._split_into_slides(raw_text)
        print(f"  [ParserAgent] Found {len(slides)} slides/sections")

        # ── Step 3: spaCy extracts noun phrases ──────────
        print("\n  [ParserAgent] REASON Step 2: Running spaCy NLP...")
        candidates = self._extract_noun_phrases(slides)
        print(f"  [ParserAgent] Raw noun phrases found: {len(candidates)}")

        # ── Step 4: clean and rank ────────────────────────
        clean_concepts = self._clean_candidates(candidates)
        print(f"  [ParserAgent] After cleaning: {len(clean_concepts)} concepts")
        print(f"  [ParserAgent] Candidates → {clean_concepts}")
        self.blackboard.log_thinking(
            "ParserAgent",
            f"NLP found {len(clean_concepts)} candidate concepts: {clean_concepts[:5]}..."
        )

        # ── Step 5: LLM builds dependency graph ──────────
        print("\n  [ParserAgent] REASON Step 3: Asking LLM to build concept graph...")
        concept_graph = self._build_concept_graph(clean_concepts, raw_text[:3000])
        print(f"  [ParserAgent] Concept graph built with {len(concept_graph)} concepts")

        return concept_graph

    # ── ACT ──────────────────────────────────────────────────
    def act(self, decision: list):
        """Write the concept graph and set the first concept."""
        self.blackboard.write("concept_graph",   decision)
        self.blackboard.write("current_concept", decision[0]["concept"])

        print("\n  [ParserAgent] ACT → Writing to blackboard:")
        print(f"  {'─'*40}")
        for i, c in enumerate(decision):
            deps = c.get("depends_on", [])
            dep_str = f" (needs: {deps})" if deps else " (no prerequisites)"
            print(f"  {i+1}. [{c['difficulty']}/5] {c['concept']}{dep_str}")
        print(f"  {'─'*40}")
        print(f"  [ParserAgent] Starting concept: {decision[0]['concept']}")

        self.blackboard.log_thinking(
            "ParserAgent",
            f"Concept graph ready. Found {len(decision)} concepts. "
            f"Starting with: '{decision[0]['concept']}' (easiest first)."
        )

    # ── PRIVATE HELPERS ───────────────────────────────────────

    def _extract_text(self, path: str) -> str:
        """Extract raw text from any supported file format."""
        ext = path.lower().split(".")[-1]

        if ext in ["ppt", "pptx"]:
            return self._extract_from_pptx(path)
        elif ext == "pdf":
            return self._extract_from_pdf(path)
        elif ext in ["docx", "doc"]:
            return self._extract_from_docx(path)
        elif ext == "txt":
            with open(path, "r", encoding="utf-8") as f:
                return f.read()
        else:
            raise ValueError(f"Unsupported file format: .{ext}")

    def _extract_from_pptx(self, path: str) -> str:
        """Extract text from PowerPoint using python-pptx."""
        try:
            from pptx import Presentation
            prs = Presentation(path)
            slides_text = []
            for slide in prs.slides:
                slide_parts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_parts.append(shape.text.strip())
                if slide_parts:
                    slides_text.append("\f".join(slide_parts))
            return "\f".join(slides_text)
        except ImportError:
            # fallback to markitdown
            return self._extract_with_markitdown(path)

    def _extract_from_pdf(self, path: str) -> str:
        """Extract text from PDF using PyMuPDF."""
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(path)
            pages = []
            for page in doc:
                pages.append(page.get_text())
            return "\f".join(pages)
        except ImportError:
            return self._extract_with_markitdown(path)

    def _extract_from_docx(self, path: str) -> str:
        """Extract text from Word document."""
        try:
            from docx import Document
            doc = Document(path)
            return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
        except ImportError:
            return self._extract_with_markitdown(path)

    def _extract_with_markitdown(self, path: str) -> str:
        """Universal fallback using markitdown library."""
        try:
            from markitdown import MarkItDown
            md = MarkItDown()
            result = md.convert(path)
            return result.text_content
        except ImportError:
            raise ImportError(
                "Install markitdown: pip install markitdown\n"
                "Or install specific parsers:\n"
                "  pip install python-pptx  (for .pptx)\n"
                "  pip install pymupdf      (for .pdf)\n"
                "  pip install python-docx  (for .docx)"
            )

    def _split_into_slides(self, raw_text: str) -> list:
        """Split text into slide/page sections."""
        # form feed character \f separates slides in most extractors
        slides = raw_text.split("\f")
        # also split on double newlines as backup
        if len(slides) < 3:
            slides = re.split(r"\n{3,}", raw_text)
        # filter out very short sections (slide numbers, blank slides)
        slides = [s.strip() for s in slides if len(s.strip()) > 30]
        return slides

    def _extract_noun_phrases(self, slides: list) -> list:
        """Use spaCy to extract meaningful noun phrases from all slides."""
        all_phrases = []
        for slide in slides:
            doc = self.nlp(slide[:1000])  # limit per slide for speed
            for chunk in doc.noun_chunks:
                text = chunk.text.lower().strip()
                # keep only multi-word phrases (more likely to be concepts)
                if len(text.split()) >= 2 and len(text) > 5:
                    all_phrases.append(text)
        return all_phrases

    def _clean_candidates(self, phrases: list) -> list:
        """
        Clean, deduplicate, and rank candidate concept phrases.
        Keep only frequent ones — frequency = importance in the lecture.
        """
        # count frequency
        freq = Counter(phrases)

        # remove generic stop-phrases that are never real concepts
        stopwords = {
            "the agent", "the system", "the user", "a number",
            "the process", "this part", "the problem", "the result",
            "the following", "the same", "the next", "the first",
            "the world", "the environment", "the set"
        }

        # keep if appears ≥ 2 times and not a stopword
        clean = [
            phrase for phrase, count in freq.most_common(30)
            if count >= 1 and phrase not in stopwords
        ]

        # remove substrings (e.g., if "agent" and "reactive agent" both appear, keep "reactive agent")
        final = []
        for phrase in clean:
            is_substring = any(
                phrase != other and phrase in other
                for other in clean
            )
            if not is_substring:
                final.append(phrase)

        return final[:20]  # max 20 candidates

    def _build_concept_graph(self, candidates: list, lecture_context: str) -> list:
        """
        Ask LLM to order candidates by difficulty and build dependency links.
        This is the only NLP step that requires semantic understanding.
        """
        prompt = f"""You are a curriculum designer for a university AI course.

Here is context from the lecture:
\"\"\"{lecture_context}\"\"\"

Here are candidate concepts extracted from the lecture:
{candidates}

Your task:
1. Select the most important 6-10 concepts from the candidates
2. Order them from easiest to hardest
3. Identify which concepts depend on others

Return ONLY a valid JSON array. Each item must have:
- "concept": clean, properly capitalized concept name (string)
- "difficulty": integer 1 (easiest) to 5 (hardest)  
- "depends_on": list of concept names that must be understood first (can be empty list)
- "summary": one sentence explaining this concept

Example format:
[
  {{"concept": "What is an Agent", "difficulty": 1, "depends_on": [], "summary": "An agent is a system that perceives its environment and takes actions."}},
  {{"concept": "Agent Architecture", "difficulty": 2, "depends_on": ["What is an Agent"], "summary": "The internal design blueprint that determines how an agent processes information and decides actions."}}
]

Return ONLY the JSON array, no explanation, no markdown code blocks."""

        response = ollama.chat(
            model="llama3",
            messages=[{"role": "user", "content": prompt}]
        )

        raw = response["message"]["content"].strip()

        # strip markdown code blocks if LLM added them
        raw = re.sub(r"```json\s*", "", raw)
        raw = re.sub(r"```\s*", "", raw)

        try:
            graph = json.loads(raw)
            # sort by difficulty to guarantee order
            graph.sort(key=lambda x: x.get("difficulty", 1))
            return graph
        except json.JSONDecodeError as e:
            print(f"  [ParserAgent] WARNING: LLM returned invalid JSON: {e}")
            print(f"  [ParserAgent] Raw response: {raw[:200]}")
            # fallback: create simple graph from candidates
            return [
                {
                    "concept": c.title(),
                    "difficulty": min(i + 1, 5),
                    "depends_on": [],
                    "summary": f"Concept: {c}"
                }
                for i, c in enumerate(candidates[:8])
            ]
