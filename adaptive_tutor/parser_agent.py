"""
parser_agent.py
───────────────
Responsible for:
  1. Extracting raw text from uploaded lecture (.ppt/.pdf/.docx)
  2. Running spaCy NLP to find candidate concepts
  3. Calling LLM to build ordered concept dependency graph
  4. Writing concept_graph to the blackboard

KEY IMPROVEMENT over original:
  After building the concept graph, the parser extracts the most relevant
  paragraph(s) for each concept from the full lecture text and stores them
  under the blackboard key "concept_contexts" (a dict: concept → text).
  QuestionAgent and FeedbackAgent consume this for grounded, accurate questions.

Perceives:  file_path  (on blackboard)
Acts:       concept_graph, current_concept, concept_contexts  (written to blackboard)
"""

import json
import re
from collections import Counter

import spacy

from base_agent   import BaseAgent
from llm_provider import call_llm, parse_json


class ParserAgent(BaseAgent):

    def __init__(self, blackboard):
        super().__init__("ParserAgent", blackboard)

        print("  [ParserAgent] Loading spaCy model...")
        try:
            self.nlp = spacy.load("en_core_web_sm")
        except OSError:
            print("  Falling back to lightweight text heuristics.")
            self.nlp = None

    # ── PERCEIVE ─────────────────────────────────────────────
    def perceive(self):
        file_path = self.blackboard.read("file_path")
        print(f"  [ParserAgent] PERCEIVE → file path: {file_path}")
        self.blackboard.log_thinking(
            "ParserAgent",
            f"I see a lecture file at: {file_path}. I will extract its concepts."
        )
        return {"file_path": file_path}

    # ── REASON ───────────────────────────────────────────────
    def reason(self, perception: dict) -> dict:
        """
        Full pipeline:
          raw text → slides → noun phrases → clean candidates → LLM concept graph
          → per-concept context extraction
        Returns a dict with "graph" and "contexts".
        """
        path = perception["file_path"]

        # Step 1: extract raw text
        print("\n  [ParserAgent] REASON Step 1: Extracting text from file...")
        raw_text = self._extract_text(path)
        print(f"  [ParserAgent] Extracted {len(raw_text)} characters of text")
        self.blackboard.log_thinking(
            "ParserAgent",
            f"Extracted {len(raw_text)} characters. Splitting into slides..."
        )

        # Step 2: split into slides
        slides = self._split_into_slides(raw_text)
        print(f"  [ParserAgent] Found {len(slides)} slides/sections")

        # Step 3: spaCy extracts noun phrases
        print("\n  [ParserAgent] REASON Step 2: Running spaCy NLP...")
        candidates = self._extract_noun_phrases(slides)
        print(f"  [ParserAgent] Raw noun phrases found: {len(candidates)}")

        # Step 4: clean and rank
        clean_concepts = self._clean_candidates(candidates)
        print(f"  [ParserAgent] After cleaning: {len(clean_concepts)} concepts")
        print(f"  [ParserAgent] Candidates → {clean_concepts}")
        self.blackboard.log_thinking(
            "ParserAgent",
            f"NLP found {len(clean_concepts)} candidate concepts: {clean_concepts[:5]}..."
        )

        # Step 5: LLM builds dependency graph
        print("\n  [ParserAgent] REASON Step 3: Asking LLM to build concept graph...")
        concept_graph = self._build_concept_graph(clean_concepts, raw_text[:4000])
        print(f"  [ParserAgent] Concept graph built with {len(concept_graph)} concepts")

        # Step 6 (NEW): extract per-concept contexts from full lecture text
        print("\n  [ParserAgent] REASON Step 4: Extracting per-concept lecture contexts...")
        concept_contexts = self._extract_concept_contexts(concept_graph, raw_text, slides)
        print(f"  [ParserAgent] Built contexts for {len(concept_contexts)} concepts")

        return {"graph": concept_graph, "contexts": concept_contexts}

    # ── ACT ──────────────────────────────────────────────────
    def act(self, decision: dict):
        """Write the concept graph, concept contexts, and set the first concept."""
        graph    = decision["graph"]
        contexts = decision["contexts"]

        self.blackboard.write("concept_graph",    graph)
        self.blackboard.write("concept_contexts", contexts)
        self.blackboard.write("current_concept",  graph[0]["concept"])

        print("\n  [ParserAgent] ACT → Writing to blackboard:")
        print(f"  {'─'*40}")
        for i, c in enumerate(graph):
            deps    = c.get("depends_on", [])
            dep_str = f" (needs: {deps})" if deps else " (no prerequisites)"
            ctx_len = len(contexts.get(c['concept'], ''))
            print(f"  {i+1}. [{c['difficulty']}/5] {c['concept']}{dep_str}  [ctx: {ctx_len} chars]")
        print(f"  {'─'*40}")
        print(f"  [ParserAgent] Starting concept: {graph[0]['concept']}")

        self.blackboard.log_thinking(
            "ParserAgent",
            f"Concept graph ready. Found {len(graph)} concepts with per-concept contexts. "
            f"Starting with: '{graph[0]['concept']}' (easiest first)."
        )

    # ── PER-CONCEPT CONTEXT EXTRACTION (NEW) ─────────────────
    def _extract_concept_contexts(
        self, graph: list, raw_text: str, slides: list
    ) -> dict:
        """
        For each concept in the graph, find the most relevant slide(s) / paragraph(s)
        from the lecture text and store up to ~800 characters.

        Strategy (in priority order):
          1. Find slides that contain the concept name (case-insensitive).
          2. If none found, use LLM to identify the best passage (for short lectures).
          3. Fall back to the concept's own summary from the graph.
        """
        contexts = {}

        for node in graph:
            concept = node["concept"]
            summary = node.get("summary", "")

            # ── Strategy 1: keyword search in slides ──────────
            keyword      = concept.lower()
            # also try the first word if multi-word (e.g. "Bayesian Network" → "bayesian")
            keyword_alt  = keyword.split()[0] if " " in keyword else None

            matching_slides = []
            for slide in slides:
                slide_lower = slide.lower()
                if keyword in slide_lower or (keyword_alt and keyword_alt in slide_lower):
                    matching_slides.append(slide)

            if matching_slides:
                # Concatenate up to 3 most relevant slides, trimmed to 800 chars
                combined = "\n\n".join(matching_slides[:3])
                contexts[concept] = combined[:800].strip()
                continue

            # ── Strategy 2: search raw_text by paragraph ──────
            paragraphs = [p.strip() for p in re.split(r"\n{2,}", raw_text) if p.strip()]
            matching_paras = [
                p for p in paragraphs
                if keyword in p.lower() or (keyword_alt and keyword_alt in p.lower())
            ]
            if matching_paras:
                combined = "\n\n".join(matching_paras[:4])
                contexts[concept] = combined[:800].strip()
                continue

            # ── Strategy 3: fall back to summary from graph ───
            # Supplement with any surrounding text from raw_text using a wider search
            first_occurrence = raw_text.lower().find(keyword)
            if first_occurrence != -1:
                start  = max(0, first_occurrence - 200)
                end    = min(len(raw_text), first_occurrence + 600)
                snippet = raw_text[start:end].strip()
                contexts[concept] = snippet[:800]
            else:
                contexts[concept] = summary  # last resort

        return contexts

    # ── PRIVATE HELPERS ───────────────────────────────────────

    def _extract_text(self, path: str) -> str:
        ext = path.lower().split(".")[-1]
        if ext in ["ppt", "pptx"]:
            return self._extract_from_pptx(path)
        elif ext == "pdf":
            return self._extract_from_pdf(path)
        elif ext in ["docx", "doc"]:
            return self._extract_from_docx(path)
        elif ext == "txt":
            with open(path, "r", encoding="utf-8") as f:
                return f.read()
        else:
            raise ValueError(f"Unsupported file format: .{ext}")

    def _extract_from_pptx(self, path: str) -> str:
        try:
            from pptx import Presentation
            prs = Presentation(path)
            slides_text = []
            for slide in prs.slides:
                slide_parts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_parts.append(shape.text.strip())
                if slide_parts:
                    slides_text.append("\f".join(slide_parts))
            return "\f".join(slides_text)
        except ImportError:
            return self._extract_with_markitdown(path)

    def _extract_from_pdf(self, path: str) -> str:
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(path)
            pages = []
            for page in doc:
                pages.append(page.get_text())
            return "\f".join(pages)
        except ImportError:
            return self._extract_with_markitdown(path)

    def _extract_from_docx(self, path: str) -> str:
        try:
            from docx import Document
            doc = Document(path)
            return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
        except ImportError:
            return self._extract_with_markitdown(path)

    def _extract_with_markitdown(self, path: str) -> str:
        try:
            from markitdown import MarkItDown
            md = MarkItDown()
            result = md.convert(path)
            return result.text_content
        except ImportError:
            raise ImportError(
                "Install markitdown: pip install markitdown\n"
                "Or install specific parsers:\n"
                "  pip install python-pptx  (for .pptx)\n"
                "  pip install pymupdf      (for .pdf)\n"
                "  pip install python-docx  (for .docx)"
            )

    def _split_into_slides(self, raw_text: str) -> list:
        slides = raw_text.split("\f")
        if len(slides) < 3:
            slides = re.split(r"\n{3,}", raw_text)
        slides = [s.strip() for s in slides if len(s.strip()) > 30]
        return slides

    def _extract_noun_phrases(self, slides: list) -> list:
        if self.nlp is None:
            return self._extract_candidate_phrases_fallback(slides)
        all_phrases = []
        for slide in slides:
            doc = self.nlp(slide[:1000])
            for chunk in doc.noun_chunks:
                text = chunk.text.lower().strip()
                if len(text.split()) >= 2 and len(text) > 5:
                    all_phrases.append(text)
        return all_phrases

    def _extract_candidate_phrases_fallback(self, slides: list) -> list:
        phrases = []
        for slide in slides:
            for raw_line in slide.splitlines():
                line = re.sub(r"[^A-Za-z0-9\-\s]", " ", raw_line).strip()
                if not line:
                    continue
                words = [w for w in line.lower().split() if len(w) > 2]
                if 2 <= len(words) <= 6:
                    phrases.append(" ".join(words))
                for n in (2, 3, 4):
                    for i in range(len(words) - n + 1):
                        gram = words[i:i + n]
                        if all(word.isalpha() for word in gram):
                            phrases.append(" ".join(gram))
        return phrases

    def _clean_candidates(self, phrases: list) -> list:
        freq = Counter(phrases)
        stopwords = {
            "the agent", "the system", "the user", "a number",
            "the process", "this part", "the problem", "the result",
            "the following", "the same", "the next", "the first",
            "the world", "the environment", "the set"
        }
        clean = [
            phrase for phrase, count in freq.most_common(30)
            if count >= 1 and phrase not in stopwords
        ]
        final = []
        for phrase in clean:
            is_substring = any(
                phrase != other and phrase in other
                for other in clean
            )
            if not is_substring:
                final.append(phrase)
        return final[:20]

    def _build_concept_graph(self, candidates: list, lecture_context: str) -> list:
        """
        Ask the LLM to select real teachable concepts from the candidate list
        and build an ordered dependency graph with proper summaries.
        Uses more lecture context (4 000 chars) for better accuracy.
        """
        prompt = f"""You are an expert curriculum designer analyzing a university lecture.

Here are candidate topics extracted from the lecture (may include noise):
{candidates}

Here is the lecture content (first 4000 characters):
\"\"\"
{lecture_context}
\"\"\"

Your task: select 5-8 REAL, teachable concepts from this lecture.

STRICT RULES:
1. Each concept must be a proper academic topic (e.g. "Natural Language Processing", "Text Corpora", "Tokenization").
2. NEVER include: professor names, emails, course codes, "agenda", "contents", sentence fragments, or anything that is not a topic.
3. Order concepts from simplest to most complex (concept 1 has no prerequisites).
4. The "summary" field must be 2-3 sentences explaining what the concept IS, taken DIRECTLY from the lecture content above.
5. "difficulty" must be 1 (easiest) to 5 (hardest).
6. "depends_on" must be concept names that must be understood first (empty list for first concept).

Return ONLY valid JSON, no markdown:
[
  {{
    "concept": "Natural Language",
    "difficulty": 1,
    "depends_on": [],
    "summary": "A natural language is one developed by humans through natural use and communication, as opposed to artificial programming languages. It is the primary object of study in NLP."
  }},
  ...
]"""

        try:
            raw   = call_llm(prompt, max_tokens=1000)
            graph = parse_json(raw)

            forbidden_in_concept = {"dr.", "prof", "email", "@", "university", "agenda",
                                     "contents", "maryam", "noha", "nourhan", "copyright",
                                     ".edu", ".com", "csc"}
            forbidden_in_summary = {"@", ".edu", ".com", "http", "dr.", "prof",
                                     "miuegypt", "csc0275"}
            cleaned = [
                c for c in graph
                if isinstance(c.get("concept"), str)
                and isinstance(c.get("summary"), str)
                and len(c["concept"]) <= 60
                and len(c["concept"].split()) <= 7
                and not c["concept"].endswith((",", "."))
                and not any(f in c["concept"].lower() for f in forbidden_in_concept)
                and len(c["summary"]) >= 30
                and " " in c["summary"]
                and not any(f in c["summary"].lower() for f in forbidden_in_summary)
            ]

            if len(cleaned) < 3:
                raise ValueError(f"Too few valid concepts after filtering: {cleaned}")

            return cleaned

        except Exception as e:
            print(f"  [ParserAgent] LLM concept graph failed ({e}), using fallback...")
            return self._fallback_graph_from_text(lecture_context)

    def _fallback_graph_from_text(self, text: str) -> list:
        garbage_patterns = [
            "dr.", "prof", "email", "@", ".edu", ".com", ".org",
            "contents", "agenda", "lecture", "university", "copyright",
            "page", "slide", "textbook", "marks", "attendance",
            "midterm", "final exam", "csc", "miuegypt", "eng.",
        ]

        def is_garbage(s: str) -> bool:
            low = s.lower()
            return any(g in low for g in garbage_patterns)

        def looks_like_sentence(s: str) -> bool:
            if is_garbage(s):
                return False
            if len(s) < 25:
                return False
            if " " not in s:
                return False
            if not s[0].isalpha():
                s = s.lstrip("•–-– ").strip()
                if not s or not s[0].isalpha():
                    return False
            return True

        lines    = text.splitlines()
        headings = []
        seen     = set()

        for i, raw_line in enumerate(lines):
            line = raw_line.strip()
            if not line or len(line) < 5 or len(line) > 60:
                continue
            if is_garbage(line):
                continue
            if line.endswith(",") or not line[0].isupper():
                continue
            word_count = len(line.split())
            if word_count < 2 or word_count > 6:
                continue
            if line in seen:
                continue
            seen.add(line)

            summary = ""
            for j in range(i + 1, min(i + 8, len(lines))):
                candidate = lines[j].strip().lstrip("•–-– ").strip()
                if looks_like_sentence(candidate):
                    summary = candidate[:180]
                    break

            if not summary:
                summary = f"{line} is a core concept covered in this lecture."

            headings.append((line, summary))

        valid = headings[2:10] if len(headings) > 4 else headings
        return [
            {
                "concept":    h,
                "difficulty": min(i + 1, 5),
                "depends_on": [],
                "summary":    s,
            }
            for i, (h, s) in enumerate(valid)
        ]